import streamlit as st
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
import io


def calculate_goal_progress(actual, projected):
    try:
        actual = float(actual)
        projected = float(projected)
        if projected == 0:
            return "N/A"
        return f"{(actual / projected) * 100:.1f}%"
    except ValueError:
        return "Invalid"


def create_presentation(current_performance, projected_performance, table_data):
    prs = Presentation()
    prs.slide_width = Inches(13.33)  # Set slide width to 16:9 width
    prs.slide_height = Inches(7.5)  # Set slide height to 16:9 height
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)

    title_shape = slide.shapes.title
    title_shape.text = f"{current_performance} Performance and {projected_performance} Projections Data"
    title_shape.text_frame.paragraphs[0].font.size = Pt(24)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(16, 4, 90)
    title_shape.text_frame.paragraphs[0].font.name = 'PT Serif'
    title_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

    rows, cols = len(table_data), len(table_data[0])
    left = Inches(1)
    top = Inches(1.5)
    width = Inches(11.67)
    height = Inches(5.09)
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    table.columns[0].width = Inches(2)
    table.columns[3].width = Inches(2)

    for row_idx, row in enumerate(table_data):
        for col_idx, cell in enumerate(row):
            cell_obj = table.cell(row_idx, col_idx)
            text_frame = cell_obj.text_frame

            # Set vertical anchor
            text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE

            # Clear existing paragraphs and add a new one
            text_frame.clear()
            p = text_frame.add_paragraph()
            p.text = str(cell)
            p.alignment = PP_ALIGN.CENTER  # Horizontal alignment

            # Apply font settings
            font = p.font
            font.name = "Manrope"

            if row_idx == 0:
                font.bold = True
                font.size = Pt(14)
                font.color.rgb = RGBColor(255, 255, 255)
                cell_obj.fill.solid()
                cell_obj.fill.fore_color.rgb = RGBColor(
                    16 if col_idx < 3 else 2,
                    4 if col_idx < 3 else 81,
                    90 if col_idx < 3 else 57
                )
            else:
                font.size = Pt(12)
                font.color.rgb = RGBColor(0, 0, 0)
                cell_obj.fill.solid()
                cell_obj.fill.fore_color.rgb = (
                    RGBColor(255, 255, 255) if row_idx % 2 == 0 else RGBColor(211, 211, 211)
                )

    pptx_buffer = io.BytesIO()
    prs.save(pptx_buffer)
    pptx_buffer.seek(0)
    return pptx_buffer


def main():
    st.title("Performance and Projections PowerPoint Table Generator")

    row_names = [
        "LinkedIn Followers",
        "Email Subscribers",
        "Website Traffic (Sessions)",
        "New Client Inquires (Website)",
        "Website Qualified Leads ($5M Liquid/ $25M Net Worth)",
        "Press Mentions"
    ]

    current_year = "2024"
    projected_year = "2025"

    actual_column_data = ["200", "200", "200", "200"]
    MTD_column_data = ["10", "10", "10", "10"]
    projected_column_data = ["100", "100", "100", "100", "100", "100"]
    goal_progress_column = []

    for i, row_name in enumerate(row_names):
        if i >= 4:
            # Manual input rows
            actual = st.text_input(f"Enter {row_name} - {current_year} Actual", "")
            mtd = st.text_input(f"Enter {row_name} - {current_year} MTD", "")

            actual_column_data.append(actual)
            MTD_column_data.append(mtd)

        # Calculate goal progress
        if i < len(actual_column_data) and i < len(projected_column_data):
            goal_progress = calculate_goal_progress(actual_column_data[i], projected_column_data[i])
            goal_progress_column.append(goal_progress)

    # Ensure all lists have the same length
    while len(actual_column_data) < len(row_names):
        actual_column_data.append("")

    while len(MTD_column_data) < len(row_names):
        MTD_column_data.append("")

    while len(goal_progress_column) < len(row_names):
        goal_progress_column.append("")

    # Build table data
    table_data = [
        ["KPI", f"{current_year} Actual", f"{current_year} MTD", f"{projected_year} Projected", "Goal Progress"]]
    for i in range(len(row_names)):
        row = [
            row_names[i],
            actual_column_data[i],
            MTD_column_data[i],
            projected_column_data[i],
            goal_progress_column[i]
        ]
        table_data.append(row)

    # Generate PowerPoint and provide immediate download
    pptx_buffer = create_presentation(current_year, projected_year, table_data)

    st.download_button(
        label="Download PowerPoint File",
        data=pptx_buffer,
        file_name="generated_performance_projections.pptx",
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )


if __name__ == "__main__":
    main()
